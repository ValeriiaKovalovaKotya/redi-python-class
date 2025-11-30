from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.dml.color import RGBColor

class CV():
    def __init__(self, fname, lname, position, phone, mail, street, town):
        self.fname = fname
        self.lname = lname
        self.position = position
        self.phone = phone
        self.mail = mail
        self.street = street
        self.town = town

        self.prs = Presentation()
        self.prs.slide_height = Emu(Inches(10)) 
        self.prs.slide_width = Emu(Inches(7.5)) 

        blank_slide_layout = self.prs.slide_layouts[6] 
        self.slide = self.prs.slides.add_slide(blank_slide_layout)

        x, y, cx, cy = Inches(0.5), Inches(0.5), Inches(6.5), Inches(8.5)
        shape = self.slide.shapes.add_table(2, 2, x, y, cx, cy)

        self.table = shape.table
        self.table.rows[0].height = Inches(2)
        self.table.rows[1].height = Inches(6.5)

        self.table.columns[0].width = Inches(2.5) 
        self.table.columns[1].width = Inches(4.0)

        self.cell_0_0 = self.table.cell(0, 0)
        self.cell_0_1 = self.table.cell(0, 1)
        self.cell_1_0 = self.table.cell(1, 0)
        self.cell_1_1 = self.table.cell(1, 1)

        self.foto_x = x + Inches(0.1)
        self.foto_y = y + Inches(0.05)
        self.foto_w = Inches(2.3)
        self.foto_h = Inches(2)

        self.light_gray = RGBColor(230, 230, 230)
        self.dark_gray = RGBColor(80, 80, 80)
        self.white = RGBColor(255, 255, 255)
        self.gray = RGBColor(120, 120, 120)

        self.name = self.cell_0_1.text_frame
        self.sidebar = self.cell_1_0.text_frame
        self.sidebar.word_wrap = True 
        self.main = self.cell_1_1.text_frame
        self.main.word_wrap = True

        self.initialize_cv() 

    def add_value(self, part, title, font_size, boldiness, color, content_list = []):
        p_item = part.add_paragraph()
        p_item.font.size = Pt(font_size)
        p_item.font.bold = boldiness
        p_item.font.color.rgb = color
        
        if part == self.name: 
            p_item.text = title
        else:
            p_item.text = title.upper()
            part.add_paragraph().text = ""

            for content in content_list:
                p_content = part.add_paragraph()
                p_content.text = content
                p_content.font.size = Pt(14)
                p_content.font.color.rgb = color
            part.add_paragraph().text = ""
            
    def initialize_cv(self):
        sidebar_values = {"Contact": [self.phone, self.mail, self.street, self.town], "Skills": ["Skill1", "Skill2", "Skill3", "Skill4"], "Tongues": ["Tongue1", "Tongue2"]}
        main_values = {"Experience": ["Date", "Job_title", "Firm", "\n", "Date", "Job_title", "Firm"], "Education": ["Date", "School", "\n", "Date", "School"]}

        for i in range(2): 
            for j, cell in enumerate(self.table.rows[i].cells):
                cell.fill.solid()
                if i == 0:
                    cell.fill.fore_color.rgb = self.light_gray
                    if j == 0:
                        self.slide.shapes.add_picture('./Foto.png', self.foto_x, self.foto_y, self.foto_w, self.foto_h)
                    else:
                        self.add_value(self.name, self.fname, 24, True, self.dark_gray)
                        self.add_value(self.name, self.lname, 36, True, self.dark_gray)
                        self.add_value(self.name, self.position, 14, False, self.gray)
                elif i == 1:
                    if j == 0:
                        cell.fill.fore_color.rgb = self.dark_gray
                        for section_title, content_list in sidebar_values.items():
                            self.add_value(self.sidebar, section_title, 16, True, self.white, content_list)
                    else:
                        cell.fill.fore_color.rgb = self.white
                        for section_title, content_list in main_values.items():
                            self.add_value(self.main, section_title, 18, True, self.dark_gray, content_list)

    def save(self, filename='CV.pptx'):
        self.prs.save(filename)

my_cv = CV("Valeriia", "Kovalova", "Magition", "0123456789", "vk@gmail.de", "Street 1", "Nicetown")
my_cv.save()